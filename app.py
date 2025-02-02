import streamlit as st
import pandas as pd
import math
import os
import uuid
import random
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, UnidentifiedImageError

############################
# 1) 월별 테마 + 무작위 메시지 목록
############################
MONTHLY_THEMES = {
    1: {
        'background': 'FFE6E6',
        'title_color': 'FF3366',
        'messages': [
            '부럽다 부러워 1월의 해피버쓰?!데이',
            '1월에 태어난 당신, 최고예요!',
            '새해의 시작과 함께 태어난 1월 베이비!',
            '와우! 1월 생일이네요',
            '1월엔 눈처럼 새하얀 생일파티'
        ],
        'sub_messages': [
            '생일자 여러분 생일 당일 오후 12시 30분 퇴근하세요',
            '행복 가득한 1월 되세요',
            '1월의 주인공들, 축하합니다!',
            '오늘은 당신이 주인공',
            '파이팅 넘치는 1월 생일!'
        ]
    },
    2: {
        'background': 'E6EEFF',
        'title_color': '3366FF',
        'messages': [
            '2월의 러블리 벌스데이',
            '사랑 가득 2월, 당신의 생일',
            '2월 봄기운과 함께 찾아온 생일',
            '2월에 태어난 당신, 로맨틱 최고',
            '2월이라 2배로 축하해요!'
        ],
        'sub_messages': [
            '2월 생일자 여러분 축하합니다!',
            '따스한 봄을 기다리며',
            '생일 당일 12:30 퇴근, 잊지 마세요',
            '오늘의 주인공, 축하해요!',
            '생일파티는 2월 스타일'
        ]
    },
    # ... 생략 없이 3~12월도 모두 같은 구조 ...
    3: {
        'background': 'E6FFE6',
        'title_color': '33CC33',
        'messages': [
            '3월의 새싹같은 생일',
            '봄바람 타고 온 3월의 주인공',
            '새 시작, 3월 생일 축하!',
            '3월 생일, 봄꽃처럼 피어나는 날',
            '화사한 3월, 생일도 화사하게!'
        ],
        'sub_messages': [
            '봄날에 태어난 여러분!',
            '3월 생일, 설레는 봄의 기운',
            '생일 당일 오후 12:30 퇴근 어때요?',
            '프레시한 3월 생일파티',
            '꽃향기 가득한 3월'
        ]
    },
    4: {
        'background': 'FFFBE6',
        'title_color': 'CC9900',
        'messages': [
            '4월의 화사한 생일!',
            '꽃 피는 4월, 생일 축하',
            '4월엔 봄바람과 함께 생일',
            '4월의 따스한 햇살과 함께',
            '4월이라 더욱 화사한 당신!'
        ],
        'sub_messages': [
            '꽃처럼 피어난 4월의 주인공',
            '생일 당일은 즐겁게!',
            '봄 날씨만큼 기분 좋은 생일',
            '4월, 꽃길 걷는 생일파티',
            '축하 폭주! 4월의 기쁨'
        ]
    },
    5: {
        'background': 'F0F0F0',
        'title_color': '333333',
        'messages': [
            '5월의 가정의 달 탄생!',
            '행복한 5월, 생일도 행복!',
            '가정의 달에 태어난 축복',
            '5월이라 가족과 함께 두 배로',
            '푸르른 5월, 생일 축하합니다'
        ],
        'sub_messages': [
            '가족 같은 회사에서 함께해요',
            '5월의 주인공, 축하해요!',
            '싱그러운 5월의 생일',
            '생일 당일 오후 12:30 퇴근합시다',
            '웃음 가득한 5월'
        ]
    },
    6: {
        'background': 'E6FFFF',
        'title_color': '00CCCC',
        'messages': [
            '6월의 시원한 생일이!',
            '초여름 바람과 함께하는 생일',
            '6월 생일, 여름의 시작!',
            '생일도 시원하게, 6월',
            '햇살 가득 6월 생일'
        ],
        'sub_messages': [
            '여름을 시원하게 만들어줄 당신',
            '6월 생일자, 당신이 최고',
            '얼음 가득한 파티 준비 OK?',
            '오늘은 당신이 주인공!',
            '시원한 생일파티 즐기세요'
        ]
    },
    7: {
        'background': 'FFF0F5',
        'title_color': 'FF0066',
        'messages': [
            '7월의 뜨거운 생일',
            '한여름 열정 가득, 7월!',
            '7월에 태어난 태양 같은 당신',
            '무더위보다 뜨거운 생일파티',
            '7월엔 열정 파워 업'
        ],
        'sub_messages': [
            '열정 가득 7월의 주인공',
            '더운 날씨도 잊을 생일 축하',
            '생일 당일 반차 어때요?',
            '7월엔 시원한 음료와 함께',
            '뜨거운 마음으로 축하합니다'
        ]
    },
    8: {
        'background': 'F5F5DC',
        'title_color': '996600',
        'messages': [
            '8월의 태양처럼! 생일',
            '뜨거운 여름, 뜨거운 생일',
            '8월에 태어난 열정맨/열정우먼',
            '태양보다 더 빛나는 8월 생일',
            '8월 무더위 속 시원한 파티'
        ],
        'sub_messages': [
            '한여름 태양보다 뜨거운 축하',
            '8월의 주인공, 당신!',
            '생일 당일 오전 근무만?',
            '시원한 바캉스와 생일파티',
            '파워풀한 8월 생일!'
        ]
    },
    9: {
        'background': 'F0FFF0',
        'title_color': '009966',
        'messages': [
            '9월의 풍성한 생일',
            '가을 문턱, 9월 생일',
            '풍요로운 9월, 생일도 풍성',
            '9월에 태어난 당신, 수확의 기쁨',
            '가을바람과 함께 9월 생일'
        ],
        'sub_messages': [
            '가을처럼 풍요로운 9월',
            '감성 가득 9월 생일',
            '생일 당일 오후 12:30 퇴근~',
            '9월엔 낙엽과 함께 축하해요',
            '오늘만은 당신이 주인공!'
        ]
    },
    10: {
        'background': 'FFFACD',
        'title_color': 'CC6600',
        'messages': [
            '10월의 청명한 생일',
            '하늘이 높은 10월, 생일 축하',
            '가을 정취 가득 10월',
            '맑고 높은 하늘처럼 빛나는 당신',
            '10월엔 단풍과 함께 생일'
        ],
        'sub_messages': [
            '맑고 높은 하늘처럼 빛나는',
            '가을 속에 피어난 생일',
            '생일 당일 오후 반차!',
            '10월 생일, 낭만 가득',
            '가을 하늘만큼 푸른 축하'
        ]
    },
    11: {
        'background': 'F5F5F5',
        'title_color': '666666',
        'messages': [
            '11월의 감사하는 생일',
            '늦가을의 아름다움, 11월',
            '11월에 태어난 당신, 고마워요',
            '가을의 끝, 11월 생일',
            '찬바람에도 따뜻한 11월'
        ],
        'sub_messages': [
            '가을의 끝, 감사의 마음과 함께',
            '생일 당일 반차 신청!',
            '11월 생일자, 축하합니다',
            '포근한 담요 같은 생일',
            '늦가을 감성 파티'
        ]
    },
    12: {
        'background': 'FFEFFC',
        'title_color': 'FF33CC',
        'messages': [
            '12월의 따뜻한 생일',
            '연말에 더 반짝이는 당신!',
            '12월, 크리스마스와 함께',
            '한 해의 마무리를 장식하는 생일',
            '겨울 분위기 가득한 12월'
        ],
        'sub_messages': [
            '연말에 더 빛나는 생일',
            '눈 내리는 12월의 파티',
            '생일 당일 오후 12:30 퇴근 가능?',
            '올 한 해의 주인공은 당신',
            '아듀 12월, 생일 축하!'
        ]
    }
}

############################
# n명일 때 행별 배치 예시 (단순)
############################
def get_rows_for_n(n: int):
    if n == 1:
        return [1]
    elif n == 2:
        return [2]
    elif n == 3:
        return [3]
    elif n == 4:
        return [4]
    elif n == 5:
        return [5]
    elif n == 6:
        return [6]
    elif n == 7:
        return [4,3]
    elif n == 8:
        return [4,4]
    elif n == 9:
        return [5,4]
    elif n == 10:
        return [5,5]
    elif n == 11:
        return [6,5]
    elif n == 12:
        return [6,6]
    elif n == 13:
        return [7,6]
    else:
        return [n]

############################
# PPT 생성 클래스
############################
class BirthdaySlideGenerator:
    def __init__(self, month: int = 1):
        self.prs = Presentation()
        self.prs.slide_width = int(33.87 * 360000)
        self.prs.slide_height = int(19.05 * 360000)
        
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        self.SLIDE_WIDTH = self.prs.slide_width
        self.SLIDE_HEIGHT = self.prs.slide_height
        
        self.month = month
        self.theme = MONTHLY_THEMES.get(month, {
            'background': 'FFFFFF',
            'title_color': '000000',
            'messages': ['기본 메시지'],
            'sub_messages': ['기본 서브']
        })
        
        # 이 달의 메시지 중 무작위
        self.main_message = random.choice(self.theme['messages'])
        self.sub_message = random.choice(self.theme['sub_messages'])

    def create_circle_image(self, image_path, output_path, size=(300,300)):
        img = Image.open(image_path).convert("RGBA")
        img = img.resize(size, Image.Resampling.LANCZOS)
        
        mask = Image.new('L', size, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, size[0], size[1]), fill=255)
        
        output = Image.new('RGBA', size, (0, 0, 0, 0))
        output.paste(img, (0,0), mask=mask)
        output.save(output_path, 'PNG')
        return output_path

    def add_decorations(self):
        slide_width_in = self.SLIDE_WIDTH / 914400.0
        shapes = [MSO_SHAPE.PENTAGON, MSO_SHAPE.OVAL, MSO_SHAPE.DIAMOND]
        colors = ['FFD700', 'FF69B4', 'FF6B6B']
        
        # 왼쪽 상단
        for i in range(3):
            shape = self.slide.shapes.add_shape(
                shapes[i % len(shapes)],
                Inches(0.5 + i * 1.5),
                Inches(0.3),
                Inches(0.4),
                Inches(0.4)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(colors[i % len(colors)])
            shape.line.width = Pt(0)
        
        # 오른쪽 상단
        for i in range(3):
            x_pos = slide_width_in - (0.5 + i * 1.5)
            shape = self.slide.shapes.add_shape(
                shapes[i % len(shapes)],
                Inches(x_pos),
                Inches(0.3),
                Inches(0.4),
                Inches(0.4)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(colors[i % len(colors)])
            shape.line.width = Pt(0)

    def add_profile(self, info, left_emu, top_emu, photo_size=1.5):
        """
        사진 중심에 텍스트 박스를 가로로 맞추고,
        세로 간격을 최소화하여 슬라이드 밖으로 안 나가도록 조정.
        
        날짜 박스: 0.567 x 0.201 in
        텍스트 박스: 1.295 x 0.472 in
        """
        # 1) 원형 사진
        photo_h = Inches(photo_size)
        valid_photo = False
        if info.get('image_path') and os.path.exists(info.get('image_path')):
            try:
                tmp = f"temp_circle_{uuid.uuid4().hex}.png"
                self.create_circle_image(info['image_path'], tmp, (300,300))
                self.slide.shapes.add_picture(tmp, left_emu, top_emu, height=photo_h)
                os.remove(tmp)
                valid_photo = True
            except (UnidentifiedImageError, OSError):
                pass
        
        # No Photo 처리
        if not valid_photo:
            shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, left_emu, top_emu, photo_h, photo_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(200,200,200)
            shape.line.width = Pt(0)
            tf = shape.text_frame
            tf.text = "No Photo"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            for run in tf.paragraphs[0].runs:
                run.font.size = Pt(10)
                run.font.name = "Pretendard"
        
        # 2) 날짜 박스 (사진 아래 80%)
        birth_str = f"{info.get('birth_month','')}/{info.get('birth_day','')}"
        date_box_w = Inches(0.567)
        date_box_h = Inches(0.201)
        
        # 날짜 박스의 left는 사진 left
        box_left = left_emu
        box_top = top_emu + photo_h * 0.8
        
        date_box = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            box_left, box_top,
            date_box_w, date_box_h
        )
        date_box.fill.solid()
        date_box.fill.fore_color.rgb = RGBColor(255, 200, 0)
        date_box.line.width = Pt(1)
        date_box.line.color.rgb = RGBColor(255,255,255)
        
        tf_box = date_box.text_frame
        tf_box.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf_box.paragraphs[0].text = birth_str
        for run in tf_box.paragraphs[0].runs:
            run.font.name = "Pretendard"
            run.font.size = Pt(11)
            run.font.bold = True
        
        # 3) 텍스트 박스 (가로 1.295 in, 세로 0.472 in)
        #    => 사진 원형의 중앙에 맞추어, 사진 아래 minimal gap
        text_w_in = 1.295
        text_h_in = 0.472
        
        # 사진 중심 x = left_emu + photo_h / 2
        # 텍스트 박스 left = 그 중심 - text_w_in/2
        text_center_x = left_emu + photo_h/2
        text_left = text_center_x - Inches(text_w_in/2.0)
        
        # 사진 아래: top_emu + photo_h + 작게 0.01인치만 띄움
        text_top = top_emu + photo_h + Inches(0.01)
        
        box_text = self.slide.shapes.add_textbox(
            text_left, text_top,
            Inches(text_w_in), Inches(text_h_in)
        )
        tf_txt = box_text.text_frame
        tf_txt.word_wrap = True
        
        # 부서(파란색)
        dept = info.get('department','')
        eng_name = info.get('eng_name','')
        kor_name = info.get('name','')
        
        p_dept = tf_txt.add_paragraph()
        p_dept.alignment = PP_ALIGN.CENTER
        p_dept.text = dept
        for run in p_dept.runs:
            run.font.name = "Pretendard"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0,128,255)
        
        # 이름
        p_name = tf_txt.add_paragraph()
        p_name.alignment = PP_ALIGN.CENTER
        p_name.text = f"{eng_name} ({kor_name})"
        for run in p_name.runs:
            run.font.name = "Pretendard"
            run.font.size = Pt(11)

    def create_layout(self, people_data):
        n = len(people_data)
        if n == 0:
            return
        
        # 배경
        bg = self.slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor.from_string(self.theme['background'])
        
        # 제목
        title_box = self.slide.shapes.add_textbox(
            (self.SLIDE_WIDTH - Inches(12)) / 2,
            Inches(0.8),
            Inches(12),
            Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.clear()
        
        run = p.add_run()
        run.text = self.main_message
        run.font.name = "EF_레베카"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(self.theme['title_color'])
        
        # 부제목
        subtitle_box = self.slide.shapes.add_textbox(
            (self.SLIDE_WIDTH - Inches(12)) / 2,
            Inches(1.8),
            Inches(12),
            Inches(0.8)
        )
        sp = subtitle_box.text_frame.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.text = self.sub_message
        for run in sp.runs:
            run.font.name = "Pretendard"
            run.font.size = Pt(14)
        
        # 상단 장식
        self.add_decorations()

        # 행 구성
        row_cols = get_rows_for_n(n)
        
        # 셀 너비/높이 조정 → 높이를 조금 줄임(2.4) so 2줄도 슬라이드 밖 안 나가도록
        cell_w_in = 2.0
        cell_h_in = 2.4  # 전에 2.8이었는데 줄임
        
        # 상단 여백 조금 줄이거나 늘릴 수 있음
        top_margin_in = 2.5
        
        slide_w_in = self.SLIDE_WIDTH / 914400.0
        slide_h_in = self.SLIDE_HEIGHT / 914400.0
        
        idx_person = 0
        for r_i, col_count in enumerate(row_cols):
            row_w_in = col_count * cell_w_in
            row_left_in = (slide_w_in - row_w_in) / 2
            row_top_in = top_margin_in + (r_i * cell_h_in)
            
            row_left_emu = row_left_in * 914400
            row_top_emu  = row_top_in * 914400
            
            for c in range(col_count):
                if idx_person >= n:
                    break
                person = people_data[idx_person]
                left_emu = row_left_emu + c * Inches(cell_w_in)
                top_emu  = row_top_emu
                
                # 인원 1명이면 사진 좀 더 크게
                pic_size = 2.0 if n == 1 else 1.5
                
                self.add_profile(person, left_emu, top_emu, photo_size=pic_size)
                idx_person += 1

    def save(self):
        buf = BytesIO()
        self.prs.save(buf)
        buf.seek(0)
        return buf


############################
def find_image_by_engname(eng_name, image_map):
    if not eng_name:
        return None
    base = eng_name.strip().lower()
    for ext in [".png",".jpg",".jpeg"]:
        candidate = base + ext
        if candidate in image_map:
            return image_map[candidate]
    return None

############################
def main():
    st.title("AI 생일자 PPT 자동생성기")
    
    excel_file = st.file_uploader("엑셀 업로드 (xlsx, xls)", type=["xlsx","xls"])
    image_files = st.file_uploader("이미지 파일 (영문이름.png/jpg...)",
                                   type=["png","jpg","jpeg"], 
                                   accept_multiple_files=True)
    
    selected_month = st.number_input("월", 1, 12, 1)
    
    if st.button("PPT 생성"):
        if not excel_file:
            st.warning("엑셀 파일부터 올려주세요.")
            return
        
        temp_folder = f"temp_{uuid.uuid4().hex}"
        os.makedirs(temp_folder, exist_ok=True)
        
        image_map = {}
        for f in image_files:
            path = os.path.join(temp_folder, f.name)
            with open(path, "wb") as out:
                out.write(f.read())
            image_map[f.name.lower()] = path
        
        df = pd.read_excel(excel_file)  # [name, eng_name, department, birth_month, birth_day]
        df = df[df['birth_month'] == selected_month]
        if df.empty:
            st.error(f"{selected_month}월 생일자가 없습니다.")
            return
        
        people_data = []
        for _, row in df.iterrows():
            e_name = str(row.get('eng_name','')).strip()
            img_path = find_image_by_engname(e_name, image_map)
            info = {
                'name': row.get('name',''),
                'eng_name': e_name,
                'department': row.get('department',''),
                'birth_month': row.get('birth_month',''),
                'birth_day': row.get('birth_day',''),
                'image_path': img_path
            }
            people_data.append(info)
        
        if not people_data:
            st.error("생일자 데이터는 있으나, 이미지 매칭이 안 되었습니다.")
            return
        
        generator = BirthdaySlideGenerator(month=selected_month)
        generator.create_layout(people_data)
        ppt_bytes = generator.save()
        
        st.success(f"{selected_month}월 PPT 생성 완료 (총 {len(people_data)}명)")
        st.download_button(
            label="결과 PPT 다운로드",
            data=ppt_bytes,
            file_name="birthday_slide_result.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        # (선택) 임시 폴더 삭제
        # import shutil
        # shutil.rmtree(temp_folder, ignore_errors=True)

if __name__ == "__main__":
    main()
